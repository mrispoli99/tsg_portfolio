"""
export_ppt.py
=============
Generates a TSG Consumer portfolio PowerPoint deck from live SQL data.
Uses python-pptx with TSG brand colors and Arial font.

Slides generated:
  1. Cover
  2. Portfolio Overview (KPI cards + flags)
  3. Portfolio Revenue & EBITDA Trend (chart)
  4. Fund Summary Table
  5. Flags & Alerts Summary
  6+. One slide per company (KPI strip + charts)
  Last. Appendix title
"""

import io
import base64
from datetime import datetime

import pandas as pd
import plotly.graph_objects as go
import plotly.io as pio
from plotly.subplots import make_subplots
# pptx imported lazily inside page_export() to avoid
# crashing the whole app if python-pptx is not installed
import streamlit as st

from db import (
    load_portfolio_overview, load_fund_summary, load_flags,
    load_ltm_snapshot, load_quarterly,
    format_millions, format_multiple, format_pct, flag_color
)

# ---------------------------------------------------------------------------
# Brand colors (no # prefix for pptx)
# ---------------------------------------------------------------------------
C_NAVY      = RGBColor(0x07, 0x17, 0x33)
C_SLATE     = RGBColor(0x3F, 0x66, 0x80)
C_SKY       = RGBColor(0xA8, 0xCF, 0xDE)
C_XANTHOUS  = RGBColor(0xF3, 0xB5, 0x1F)
C_CELADON   = RGBColor(0x85, 0xD7, 0xB0)
C_SEA_GREEN = RGBColor(0x06, 0x86, 0x5C)
C_RED       = RGBColor(0xC0, 0x39, 0x2B)
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT     = RGBColor(0xF4, 0xF6, 0xF9)
C_BORDER    = RGBColor(0xE0, 0xE4, 0xEA)

# Hex strings for Plotly charts (no # prefix)
P_NAVY     = "071733"
P_SLATE    = "3F6680"
P_SKY      = "A8CFDE"
P_XANTHOUS = "F3B51F"
P_GREEN    = "06865C"
P_RED      = "C0392B"
P_BORDER   = "E0E4EA"

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def rgb_to_hex(color: RGBColor) -> str:
    return f"{color[0]:02X}{color[1]:02X}{color[2]:02X}"


def add_text(slide, text, x, y, w, h,
             size=11, bold=False, italic=False,
             color=C_NAVY, align=PP_ALIGN.LEFT,
             wrap=True, font="Arial"):
    txBox = slide.shapes.add_textbox(
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text)
    run.font.name  = font
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_rect(slide, x, y, w, h, fill_color=C_NAVY, line_color=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(0.5)
    else:
        shape.line.fill.background()
    return shape


def plotly_to_image(fig, width=900, height=320) -> bytes:
    """Convert Plotly figure to PNG bytes."""
    return pio.to_image(fig, format="png", width=width, height=height, scale=2)


def add_chart_image(slide, fig, x, y, w, h, width_px=900, height_px=300):
    """Render a Plotly chart and embed as image in slide."""
    img_bytes = plotly_to_image(fig, width=width_px, height=height_px)
    img_stream = io.BytesIO(img_bytes)
    slide.shapes.add_picture(img_stream, Inches(x), Inches(y), Inches(w), Inches(h))


def slide_header(slide, title: str, subtitle: str = "TSG Consumer Partners"):
    """Add the standard navy header bar to a slide."""
    add_rect(slide, 0, 0, 13.33, 0.65, fill_color=C_NAVY)
    # TSG CONSUMER logo text
    add_text(slide, "TSG", 0.2, 0.08, 1.2, 0.5,
             size=18, bold=True, color=C_WHITE)
    add_text(slide, "CONSUMER", 0.72, 0.08, 2.5, 0.5,
             size=18, bold=False, color=C_SKY)
    # Page title right-aligned
    add_text(slide, title, 9.5, 0.1, 3.6, 0.45,
             size=12, bold=False, color=C_SKY, align=PP_ALIGN.RIGHT)


def kpi_box(slide, value, label, x, y, w=2.0, h=1.0,
            value_color=C_NAVY, bg_color=C_WHITE):
    """Draw a KPI card on the slide."""
    add_rect(slide, x, y, w, h, fill_color=bg_color, line_color=C_BORDER)
    add_text(slide, value, x + 0.1, y + 0.08, w - 0.2, 0.5,
             size=20, bold=True, color=value_color, align=PP_ALIGN.CENTER)
    add_text(slide, label, x + 0.05, y + 0.6, w - 0.1, 0.35,
             size=8, bold=False, color=C_SLATE, align=PP_ALIGN.CENTER)


def flag_color_rgb(flag: str) -> RGBColor:
    return {"Red": C_RED, "Yellow": C_XANTHOUS, "Green": C_SEA_GREEN}.get(flag, C_SLATE)


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def slide_cover(prs, as_of_date: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = C_NAVY

    # Large TSG wordmark
    add_text(slide, "TSG", 1.0, 1.5, 5.0, 2.0,
             size=72, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)
    add_text(slide, "CONSUMER", 3.2, 1.5, 9.0, 2.0,
             size=72, bold=False, color=C_SKY, align=PP_ALIGN.LEFT)

    add_text(slide, "Portfolio Analytics", 1.0, 3.6, 8.0, 0.6,
             size=24, bold=False, color=C_WHITE)
    add_text(slide, f"As of {as_of_date}", 1.0, 4.3, 8.0, 0.5,
             size=16, bold=False, color=C_SKY)

    # Decorative line
    line = slide.shapes.add_shape(1, Inches(1.0), Inches(4.0), Inches(6.0), Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = C_XANTHOUS
    line.line.fill.background()


def slide_portfolio_overview(prs, overview, flags_df):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = C_LIGHT
    slide_header(slide, "Portfolio Overview")

    # KPI cards row
    kpis = [
        (format_millions(overview.get("total_tev")),           "TOTAL TEV"),
        (format_millions(overview.get("total_ltm_revenue")),    "LTM PORTFOLIO REV"),
        (format_millions(overview.get("total_ltm_adj_ebitda")), "LTM PORTFOLIO EBITDA"),
        (format_multiple(overview.get("avg_net_leverage_wtd")), "AVG NET LEVERAGE"),
        (format_pct(overview.get("avg_adj_ebitda_margin")),     "AVG EBITDA MARGIN"),
        (f"{int(overview.get('red_flag_count',0))} Red  "
         f"{int(overview.get('yellow_flag_count',0))} Yellow  "
         f"{int(overview.get('green_flag_count',0))} Green",   "KPI FLAGS"),
    ]
    x_start = 0.3
    for i, (val, label) in enumerate(kpis):
        kpi_box(slide, val, label, x_start + i * 2.15, 0.75, w=2.0, h=1.0)

    # Alert banner if red flags exist
    red_companies = flags_df[flags_df["overall_flag"] == "Red"]["company_name"].tolist()
    if red_companies:
        add_rect(slide, 0.3, 1.85, 12.7, 0.45, fill_color=RGBColor(0xFD, 0xEC, 0xEA),
                 line_color=C_RED)
        alert_text = f"🚨 Red Flags: {', '.join(red_companies)}"
        add_text(slide, alert_text, 0.4, 1.88, 12.5, 0.38,
                 size=9, bold=True, color=C_RED)

    # Flags table
    add_text(slide, "KPI Flag Summary", 0.3, 2.4, 8.0, 0.35,
             size=11, bold=True, color=C_SLATE)

    table_data = [["Company", "Sector", "LTM Revenue", "EBITDA Margin",
                   "Net Leverage", "Rev Growth", "Flag"]]
    for _, row in flags_df.sort_values("overall_flag",
                                        key=lambda x: x.map({"Red":0,"Yellow":1,"Green":2})
                                       ).head(12).iterrows():
        table_data.append([
            row.get("company_name", ""),
            "",
            format_millions(row.get("ltm_revenue")),
            format_pct(row.get("ltm_adj_ebitda_margin")),
            format_multiple(row.get("net_leverage")),
            format_pct(row.get("revenue_yoy")),
            row.get("overall_flag", ""),
        ])

    tbl = slide.shapes.add_table(
        len(table_data), 7,
        Inches(0.3), Inches(2.75),
        Inches(12.7), Inches(4.4)
    ).table

    col_widths = [2.8, 0, 1.4, 1.4, 1.4, 1.2, 1.0]
    for i, w in enumerate(col_widths):
        if w > 0:
            tbl.columns[i].width = Inches(w)

    for r_idx, row_data in enumerate(table_data):
        for c_idx, cell_text in enumerate(row_data):
            cell = tbl.cell(r_idx, c_idx)
            cell.text = cell_text
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if c_idx > 1 else PP_ALIGN.LEFT
            run = p.runs[0] if p.runs else p.add_run()
            run.font.name = "Arial"
            run.font.size = Pt(8 if r_idx > 0 else 9)
            run.font.bold = (r_idx == 0)
            if r_idx == 0:
                run.font.color.rgb = C_WHITE
                cell.fill.solid()
                cell.fill.fore_color.rgb = C_NAVY
            elif c_idx == 6:  # Flag column
                flag_val = cell_text
                run.font.color.rgb = flag_color_rgb(flag_val)
                run.font.bold = True
            else:
                run.font.color.rgb = C_NAVY
                if r_idx % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = C_LIGHT


def slide_trend_chart(prs, quarterly_df):
    """Portfolio revenue & EBITDA trend slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = C_LIGHT
    slide_header(slide, "Portfolio Revenue & EBITDA Trend")

    agg = quarterly_df.groupby("period_label").agg(
        revenue    = ("revenue", "sum"),
        adj_ebitda = ("adj_ebitda", "sum"),
        cash_flow_date = ("cash_flow_date", "max")
    ).reset_index().sort_values("cash_flow_date").tail(12)

    fig = make_subplots(specs=[[{"secondary_y": True}]])
    fig.add_trace(go.Bar(x=agg["period_label"], y=agg["revenue"],
                          name="Revenue ($M)", marker_color=f"#{P_SLATE}",
                          opacity=0.85), secondary_y=False)
    fig.add_trace(go.Bar(x=agg["period_label"], y=agg["adj_ebitda"],
                          name="Adj. EBITDA ($M)", marker_color=f"#{P_SKY}",
                          opacity=0.85), secondary_y=False)
    margin = agg["adj_ebitda"] / agg["revenue"].replace(0, float("nan"))
    fig.add_trace(go.Scatter(x=agg["period_label"], y=margin,
                              name="EBITDA Margin %", mode="lines+markers",
                              line=dict(color=f"#{P_XANTHOUS}", width=3)),
                  secondary_y=True)
    fig.update_layout(
        height=380, plot_bgcolor="white", paper_bgcolor="white",
        margin=dict(l=20, r=20, t=10, b=40), barmode="group",
        legend=dict(orientation="h", y=-0.15, font=dict(size=10)),
        font=dict(family="Arial", color=f"#{P_NAVY}", size=11),
    )
    fig.update_yaxes(tickformat="$,.0f", gridcolor=f"#{P_BORDER}", secondary_y=False)
    fig.update_yaxes(tickformat=".0%", secondary_y=True)
    fig.update_xaxes(tickangle=-45)
    add_chart_image(slide, fig, 0.5, 0.75, 12.3, 6.4, width_px=1200, height_px=480)


def slide_fund_summary(prs, fs_df):
    """Fund summary investment table."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = C_LIGHT
    slide_header(slide, "Fund Summary — Investment Table")

    cols = ["company_name","sector","investment_date","security_type",
            "entry_tev","current_tev","gross_moi","ltm_revenue","net_leverage",
            "ltm_adj_ebitda_margin","revenue_yoy","overall_flag"]
    headers = ["Company","Sector","Entry","Security","Entry TEV",
               "Current TEV","Gross MOI","LTM Rev","Net Lev","EBITDA Mgn","Rev Growth","Flag"]

    display = fs_df[cols].copy()
    display["investment_date"]      = pd.to_datetime(display["investment_date"]).dt.strftime("%b %y")
    display["entry_tev"]            = display["entry_tev"].apply(format_millions)
    display["current_tev"]          = display["current_tev"].apply(format_millions)
    display["gross_moi"]            = display["gross_moi"].apply(format_multiple)
    display["ltm_revenue"]          = display["ltm_revenue"].apply(format_millions)
    display["net_leverage"]         = display["net_leverage"].apply(format_multiple)
    display["ltm_adj_ebitda_margin"]= display["ltm_adj_ebitda_margin"].apply(format_pct)
    display["revenue_yoy"]          = display["revenue_yoy"].apply(format_pct)

    table_data = [headers] + display.values.tolist()
    row_count  = min(len(table_data), 21)

    tbl = slide.shapes.add_table(
        row_count, len(headers),
        Inches(0.2), Inches(0.75),
        Inches(12.9), Inches(6.5)
    ).table

    col_widths = [2.0, 1.4, 0.7, 1.0, 0.9, 0.9, 0.8, 0.8, 0.7, 0.8, 0.8, 0.6]
    for i, w in enumerate(col_widths):
        tbl.columns[i].width = Inches(w)

    for r_idx in range(row_count):
        row_data = table_data[r_idx]
        for c_idx, cell_text in enumerate(row_data):
            cell = tbl.cell(r_idx, c_idx)
            cell.text = str(cell_text) if cell_text is not None else "—"
            p   = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if c_idx > 3 else PP_ALIGN.LEFT
            run = p.runs[0] if p.runs else p.add_run()
            run.font.name = "Arial"
            run.font.size = Pt(7.5 if r_idx > 0 else 8)
            run.font.bold = (r_idx == 0)
            if r_idx == 0:
                run.font.color.rgb = C_WHITE
                cell.fill.solid()
                cell.fill.fore_color.rgb = C_NAVY
            elif c_idx == len(headers) - 1:
                run.font.color.rgb = flag_color_rgb(str(cell_text))
                run.font.bold = True
            else:
                run.font.color.rgb = C_NAVY
                if r_idx % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = C_LIGHT


def slide_flags_summary(prs, flags_df):
    """Flags & alerts summary slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = C_LIGHT
    slide_header(slide, "Flags & Alerts")

    red    = flags_df[flags_df["overall_flag"] == "Red"]
    yellow = flags_df[flags_df["overall_flag"] == "Yellow"]

    # Summary counts
    for i, (count, label, color) in enumerate([
        (len(red),    "Red Flags",    C_RED),
        (len(yellow), "Yellow Flags", C_XANTHOUS),
        (len(flags_df) - len(red) - len(yellow), "Green / On Track", C_SEA_GREEN),
    ]):
        add_rect(slide, 0.3 + i * 4.3, 0.75, 4.0, 0.9, fill_color=C_WHITE, line_color=C_BORDER)
        add_text(slide, str(count), 0.3 + i * 4.3 + 0.1, 0.8, 3.8, 0.5,
                 size=28, bold=True, color=color, align=PP_ALIGN.CENTER)
        add_text(slide, label, 0.3 + i * 4.3 + 0.1, 1.3, 3.8, 0.3,
                 size=9, color=C_SLATE, align=PP_ALIGN.CENTER)

    y = 1.8
    for flag_group, color, title in [
        (red,    C_RED,      "🔴 RED — Requires Immediate Attention"),
        (yellow, C_XANTHOUS, "🟡 YELLOW — Watch List"),
    ]:
        if flag_group.empty:
            continue
        add_rect(slide, 0.3, y, 12.7, 0.3, fill_color=color, line_color=None)
        add_text(slide, title, 0.4, y + 0.03, 12.5, 0.25,
                 size=9, bold=True, color=C_WHITE)
        y += 0.35

        for _, row in flag_group.iterrows():
            issues = []
            if row.get("net_leverage_flag") == "Red":
                issues.append(f"Net Lev {format_multiple(row.get('net_leverage'))}")
            if row.get("revenue_growth_flag") == "Red":
                issues.append(f"Rev Growth {format_pct(row.get('revenue_yoy'))}")
            if row.get("ebitda_margin_flag") == "Red":
                issues.append(f"EBITDA Margin {format_pct(row.get('ltm_adj_ebitda_margin'))}")

            add_rect(slide, 0.3, y, 12.7, 0.38,
                     fill_color=C_WHITE, line_color=C_BORDER)
            add_text(slide, row.get("company_name",""), 0.4, y + 0.04, 3.0, 0.3,
                     size=9, bold=True, color=C_NAVY)
            add_text(slide, " · ".join(issues) if issues else "Multiple KPI breaches",
                     3.5, y + 0.04, 9.2, 0.3,
                     size=9, color=C_RED if flag_group is red else C_XANTHOUS)
            y += 0.42

        y += 0.1


def slide_company(prs, company_name: str, flag_row, quarterly_df):
    """One slide per portfolio company."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = C_LIGHT
    slide_header(slide, f"Company Detail — {company_name}")

    if flag_row is None:
        add_text(slide, "No data available", 0.5, 1.0, 8.0, 0.5,
                 size=12, color=C_SLATE)
        return

    # Company badge
    add_rect(slide, 0.3, 0.72, 0.6, 0.5, fill_color=C_NAVY)
    add_text(slide, company_name[:2].upper(), 0.3, 0.73, 0.6, 0.48,
             size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    overall  = flag_row.get("overall_flag", "")
    flag_clr = flag_color_rgb(overall)
    add_text(slide, f"{company_name}  ●  {overall}",
             0.95, 0.73, 7.0, 0.45,
             size=14, bold=True, color=C_NAVY)

    # KPI strip
    kpis = [
        (format_millions(flag_row.get("ltm_revenue")),      "LTM Revenue"),
        (format_millions(flag_row.get("ltm_ebitda")),        "LTM EBITDA"),
        (format_pct(flag_row.get("ltm_adj_ebitda_margin")),  "EBITDA Margin"),
        (format_multiple(flag_row.get("net_leverage")),       "Net Leverage"),
        (format_pct(flag_row.get("revenue_yoy")),             "Rev Growth YoY"),
        (format_pct(flag_row.get("ebitda_yoy")),              "EBITDA Growth"),
    ]
    for i, (val, label) in enumerate(kpis):
        kpi_box(slide, val, label, 0.3 + i * 2.15, 1.28, w=2.0, h=0.85)

    # Revenue & EBITDA chart
    comp_q = quarterly_df[quarterly_df["company_name"] == company_name].copy()
    if not comp_q.empty:
        rev = comp_q["revenue"].combine_first(comp_q["net_sales"])
        fig = make_subplots(specs=[[{"secondary_y": True}]])
        fig.add_trace(go.Bar(x=comp_q["period_label"], y=rev,
                              name="Revenue", marker_color=f"#{P_NAVY}", opacity=0.8),
                      secondary_y=False)
        fig.add_trace(go.Bar(x=comp_q["period_label"], y=comp_q["adj_ebitda"],
                              name="Adj. EBITDA", marker_color=f"#{P_SLATE}", opacity=0.8),
                      secondary_y=False)
        margin = comp_q["adj_ebitda"] / rev.replace(0, float("nan"))
        fig.add_trace(go.Scatter(x=comp_q["period_label"], y=margin,
                                  name="Margin %", mode="lines+markers",
                                  line=dict(color=f"#{P_XANTHOUS}", width=2)),
                      secondary_y=True)
        fig.update_layout(
            height=260, plot_bgcolor="white", paper_bgcolor="white",
            margin=dict(l=20, r=20, t=10, b=40), barmode="group",
            legend=dict(orientation="h", y=-0.25, font=dict(size=9)),
            font=dict(family="Arial", size=9, color=f"#{P_NAVY}"),
            showlegend=True,
        )
        fig.update_yaxes(tickformat="$,.0f", gridcolor=f"#{P_BORDER}", secondary_y=False)
        fig.update_yaxes(tickformat=".0%", secondary_y=True)
        fig.update_xaxes(tickangle=-45, nticks=8)
        add_chart_image(slide, fig, 0.3, 2.2, 8.5, 5.0, width_px=860, height_px=340)

    # Flag detail box
    flag_items = []
    for metric, flag_col, val_col, fmt in [
        ("Revenue Growth", "revenue_growth_flag",  "revenue_yoy",         format_pct),
        ("EBITDA Margin",  "ebitda_margin_flag",   "ltm_adj_ebitda_margin",format_pct),
        ("Net Leverage",   "net_leverage_flag",    "net_leverage",         format_multiple),
        ("Int. Coverage",  "interest_coverage_flag","interest_coverage",   format_multiple),
    ]:
        flag_val = flag_row.get(flag_col, "")
        val      = flag_row.get(val_col)
        flag_items.append((metric, flag_val, fmt(val) if val else "—"))

    add_rect(slide, 9.0, 2.2, 4.1, 5.0, fill_color=C_WHITE, line_color=C_BORDER)
    add_text(slide, "KPI Flags", 9.1, 2.25, 3.9, 0.3,
             size=9, bold=True, color=C_SLATE)

    for i, (metric, flag_val, val_str) in enumerate(flag_items):
        y_pos = 2.6 + i * 1.12
        clr   = flag_color_rgb(flag_val)
        add_rect(slide, 9.1, y_pos, 3.8, 0.9, fill_color=C_LIGHT, line_color=C_BORDER)
        add_text(slide, metric, 9.2, y_pos + 0.05, 3.6, 0.3,
                 size=8, color=C_SLATE)
        add_text(slide, val_str, 9.2, y_pos + 0.32, 3.6, 0.45,
                 size=18, bold=True, color=clr, align=PP_ALIGN.LEFT)
        add_text(slide, flag_val, 11.5, y_pos + 0.38, 1.3, 0.35,
                 size=9, bold=True, color=clr, align=PP_ALIGN.RIGHT)


# ---------------------------------------------------------------------------
# Main export function
# ---------------------------------------------------------------------------

def build_deck(include_companies: list = None) -> bytes:
    """
    Build the full portfolio deck and return as bytes.
    include_companies: list of company names to include (None = all)
    """
    overview  = load_portfolio_overview().iloc[0]
    fs        = load_fund_summary()
    flags     = load_flags()
    quarterly = load_quarterly()  # all companies

    as_of = str(overview.get("as_of_date", datetime.now().strftime("%B %Y")))[:7] \
        if hasattr(overview, "get") else datetime.now().strftime("%B %Y")
    as_of = datetime.now().strftime("%B %Y")

    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    # Add blank layout if needed
    if len(prs.slide_layouts) < 7:
        from pptx.util import Inches
        layout = prs.slide_layouts[0]
    # Use blank layout (index 6 in most templates)
    blank_layout_idx = min(6, len(prs.slide_layouts) - 1)

    # Monkey-patch add_slide to always use blank
    original_add = prs.slides.add_slide
    def add_blank(layout=None):
        return original_add(prs.slide_layouts[blank_layout_idx])
    prs.slides.add_slide = add_blank

    # Build slides
    slide_cover(prs, as_of)
    slide_portfolio_overview(prs, overview, flags)
    slide_trend_chart(prs, quarterly)
    slide_fund_summary(prs, fs)
    slide_flags_summary(prs, flags)

    # Company slides
    companies_to_include = include_companies or flags["company_name"].tolist()
    for company_name in companies_to_include:
        flag_row_df = flags[flags["company_name"] == company_name]
        flag_row    = flag_row_df.iloc[0] if len(flag_row_df) > 0 else None
        slide_company(prs, company_name, flag_row, quarterly)

    # Save to bytes
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()


# ---------------------------------------------------------------------------
# Streamlit export page
# ---------------------------------------------------------------------------

def page_export():
    # Lazy imports — only fail here if packages missing, not at app startup
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        from pptx.util import Inches, Pt
    except ImportError:
        st.error(
            "**python-pptx** is not installed. Add `python-pptx>=0.6.21` "
            "to `requirements.txt` and redeploy."
        )
        return
    try:
        import kaleido  # noqa — needed for plotly image export
    except ImportError:
        st.warning(
            "**kaleido** is not installed — chart images will be skipped in the export. "
            "Add `kaleido` to `requirements.txt` to enable them."
        )
    NAVY  = "#071733"
    SLATE = "#3F6680"
    SKY   = "#A8CFDE"
    BORDER = "#E0E4EA"

    st.markdown(f"""
    <div style="background:{NAVY}; padding:14px 24px; border-radius:6px; margin-bottom:20px;
                display:flex; align-items:center; justify-content:space-between;">
        <div>
            <p style="color:white; font-size:20px; font-weight:700;
                      font-family:Arial; margin:0;">
                TSG <span style="color:{SKY}; font-weight:400;">CONSUMER</span>
            </p>
            <p style="color:{SKY}; font-size:13px; font-family:Arial; margin:0;">
                TSG Consumer Partners</p>
        </div>
        <p style="color:{SKY}; font-size:12px; font-family:Arial; margin:0;">
            Export to PowerPoint</p>
    </div>
    """, unsafe_allow_html=True)

    st.markdown(f"""
    <div style="background:white; border:1px solid {BORDER}; border-radius:6px;
                padding:20px; margin-bottom:16px;">
        <p style="font-size:14px; font-weight:700; color:{NAVY}; font-family:Arial; margin:0 0 8px 0;">
            Generate Portfolio Deck</p>
        <p style="font-size:12px; color:{SLATE}; font-family:Arial; margin:0;">
            Builds a branded PowerPoint presentation with cover, portfolio overview,
            trend charts, fund summary table, flags & alerts, and individual company slides.
        </p>
    </div>
    """, unsafe_allow_html=True)

    # Options
    col_opt1, col_opt2 = st.columns(2)
    with col_opt1:
        st.markdown("**Slides to include:**")
        inc_overview = st.checkbox("Portfolio Overview",   value=True)
        inc_trend    = st.checkbox("Revenue & EBITDA Trend", value=True)
        inc_fund     = st.checkbox("Fund Summary Table",   value=True)
        inc_flags    = st.checkbox("Flags & Alerts",       value=True)
        inc_companies = st.checkbox("Individual Company Slides", value=True)

    with col_opt2:
        if inc_companies:
            st.markdown("**Select companies:**")
            flags_df  = load_flags()
            all_names = flags_df["company_name"].tolist()
            selected  = st.multiselect(
                "Companies", all_names, default=all_names,
                label_visibility="collapsed"
            )
        else:
            selected = []

    st.markdown("<br>", unsafe_allow_html=True)

    if st.button("🎯 Generate PowerPoint", type="primary", use_container_width=False):
        with st.spinner("Building deck — this takes about 30 seconds..."):
            try:
                pptx_bytes = build_deck(
                    include_companies=selected if inc_companies else []
                )
                filename = f"TSG_Consumer_Portfolio_{datetime.now().strftime('%Y%m%d')}.pptx"
                st.download_button(
                    label="⬇️ Download PowerPoint",
                    data=pptx_bytes,
                    file_name=filename,
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    use_container_width=False,
                )
                st.success(f"✅ Deck ready — {len(selected) + 5} slides generated")
            except Exception as exc:
                st.error(f"Export failed: {exc}")
                raise
